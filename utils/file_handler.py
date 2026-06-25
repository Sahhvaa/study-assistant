import pdfplumber
import docx
import pptx
import openpyxl
from pathlib import Path
import fitz  # PyMuPDF for image extraction
import os

IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp'}
ALLOWED_EXTENSIONS = {
    'pdf', 'docx', 'txt', 'pptx',
    'xlsx', 'png', 'jpg', 'jpeg', 'webp'
}

def allowed_file(filename):
    return (
        '.' in filename and
        filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS
    )

def is_image(filepath):
    return Path(filepath).suffix.lower() in IMAGE_EXTENSIONS

def extract_text(filepath):
    ext = Path(filepath).suffix.lower()
    if ext == '.pdf':
        return extract_pdf(filepath)
    elif ext == '.docx':
        return extract_docx(filepath)
    elif ext == '.txt':
        return extract_txt(filepath)
    elif ext == '.pptx':
        return extract_pptx(filepath)
    elif ext == '.xlsx':
        return extract_xlsx(filepath)
    elif ext in IMAGE_EXTENSIONS:
        return None
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def extract_images_from_pdf(filepath, output_folder='uploads/extracted'):
    """Extract all images from a PDF and save them."""
    os.makedirs(output_folder, exist_ok=True)
    image_paths = []

    try:
        doc = fitz.open(filepath)
        for page_num in range(len(doc)):
            page = doc[page_num]
            images = page.get_images(full=True)

            for img_index, img in enumerate(images):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image['image']
                ext = base_image['ext']

                # save image
                image_filename = f"page{page_num+1}_img{img_index+1}.{ext}"
                image_path = os.path.join(output_folder, image_filename)

                with open(image_path, 'wb') as f:
                    f.write(image_bytes)

                image_paths.append(image_path)

        doc.close()
    except Exception as e:
        print(f"Image extraction error: {e}")

    return image_paths

def extract_pdf(filepath):
    text = ""
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_docx(filepath):
    doc = docx.Document(filepath)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_txt(filepath):
    with open(filepath, 'r', encoding='utf-8') as f:
        return f.read()

def extract_pptx(filepath):
    prs = pptx.Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_xlsx(filepath):
    wb = openpyxl.load_workbook(filepath)
    text = ""
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text += f"\n--- Sheet: {sheet_name} ---\n"
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(
                str(cell) for cell in row if cell is not None
            )
            if row_text:
                text += row_text + "\n"
    return text